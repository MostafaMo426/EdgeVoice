import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from bs4 import BeautifulSoup

def create_presentation():
    # 1. Initialize Presentation with Widescreen (16:9) aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors
    BG_COLOR = RGBColor(13, 14, 18)      # #0d0e12
    CYAN_COLOR = RGBColor(0, 242, 254)   # #00f2fe
    WHITE_COLOR = RGBColor(255, 255, 255)
    GREY_COLOR = RGBColor(148, 163, 184) # #94a3b8

    # Load and Parse HTML
    html_path = 'presentation.html'
    if not os.path.exists(html_path):
        print(f"Error: {html_path} not found.")
        return

    with open(html_path, 'r', encoding='utf-8') as f:
        soup = BeautifulSoup(f, 'html.parser')

    slides_html = soup.find_all('div', class_='slide')

    for index, slide_div in enumerate(slides_html):
        # Create a blank slide layout
        slide_layout = prs.slide_layouts[6] # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Set Background Color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        # Extract Header (h1 or h2)
        header_tag = slide_div.find(['h1', 'h2'])
        if header_tag:
            header_text = header_tag.get_text().strip()

            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(1))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = header_text
            p.font.bold = True
            p.font.size = Pt(44)
            p.font.color.rgb = CYAN_COLOR
            p.alignment = PP_ALIGN.LEFT if "divider" not in slide_div.get('class', []) else PP_ALIGN.CENTER

        # Extract Content
        # Handle Divider Slides (Centered)
        if "divider" in slide_div.get('class', []):
            subtitle_tag = slide_div.find('p')
            if subtitle_tag:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.3), Inches(1))
                p = txBox.text_frame.paragraphs[0]
                p.text = subtitle_tag.get_text().strip()
                p.font.size = Pt(28)
                p.font.color.rgb = WHITE_COLOR
                p.alignment = PP_ALIGN.CENTER
            continue

        # Handle Content Grid / Step Lists
        content_y = Inches(1.8)

        # Look for paragraphs or lists
        paragraphs = slide_div.find_all(['p', 'li', 'h3'])
        for tag in paragraphs:
            if tag.parent.name == 'li': continue # Skip individual items if we process parent

            text = tag.get_text().strip()
            if not text or text in [header_tag.get_text().strip() if header_tag else ""]:
                continue

            txBox = slide.shapes.add_textbox(Inches(0.8), content_y, Inches(11.5), Inches(0.5))
            p = txBox.text_frame.paragraphs[0]
            p.text = text

            if tag.name == 'h3':
                p.font.size = Pt(24)
                p.font.color.rgb = CYAN_COLOR
                content_y += Inches(0.5)
            else:
                p.font.size = Pt(18)
                p.font.color.rgb = GREY_COLOR
                content_y += Inches(0.4)

        # Handle Tables
        table_tag = slide_div.find('table')
        if table_tag:
            rows = table_tag.find_all('tr')
            cols_count = len(rows[0].find_all(['th', 'td']))

            shape = slide.shapes.add_table(len(rows), cols_count, Inches(1), Inches(2.5), Inches(11), Inches(4))
            table = shape.table

            for r_idx, row in enumerate(rows):
                cells = row.find_all(['th', 'td'])
                for c_idx, cell in enumerate(cells):
                    table.cell(r_idx, c_idx).text = cell.get_text().strip()
                    # Basic Styling
                    cell_p = table.cell(r_idx, c_idx).text_frame.paragraphs[0]
                    cell_p.font.size = Pt(14)
                    cell_p.font.color.rgb = WHITE_COLOR

    prs.save('presentation.pptx')
    print("Success: presentation.pptx has been generated.")

if __name__ == "__main__":
    create_presentation()
